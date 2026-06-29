from __future__ import annotations

import json
import sys
from pathlib import Path
from uuid import uuid4

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from typer.testing import CliRunner

from learnit_study import extraction, extractors
from learnit_study.cli import app


runner = CliRunner()


def local_tmp_path() -> Path:
    path = Path(".test_tmp") / uuid4().hex
    path.mkdir(parents=True, exist_ok=False)
    return path


def make_course(root: Path, *, course_id: str = "3025533") -> Path:
    course_dir = root / f"{course_id} - Demo Course"
    section_dir = course_dir / "Lecture 1"
    (section_dir / "materials").mkdir(parents=True)
    (course_dir / "manifest.json").write_text("{}", encoding="utf-8")
    (section_dir / "section_manifest.json").write_text("{}", encoding="utf-8")
    return course_dir


def read_section_manifest(course_dir: Path) -> dict:
    return json.loads((course_dir / "Lecture 1" / "section_manifest.json").read_text(encoding="utf-8"))


def write_pdf(path: Path, text: str) -> None:
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 4 0 R >> >> "
        b"/MediaBox [0 0 612 792] /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    content = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(content))
        content.extend(f"{index} 0 obj\n".encode("ascii"))
        content.extend(obj)
        content.extend(b"\nendobj\n")
    xref_start = len(content)
    content.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    content.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        content.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    content.extend(
        f"trailer << /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_start}\n%%EOF\n".encode("ascii")
    )
    path.write_bytes(bytes(content))


def test_pdf_extraction() -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    pdf_path = course_dir / "Lecture 1" / "materials" / "lecture.pdf"
    write_pdf(pdf_path, "PDF hello")

    extraction.extract_course_text(course_dir=course_dir)

    extracted = course_dir / "Lecture 1" / "extracted_text.md"
    assert "PDF hello" in extracted.read_text(encoding="utf-8")


def test_docx_extraction() -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    doc = Document()
    doc.add_paragraph("DOCX hello")
    doc.save(course_dir / "Lecture 1" / "materials" / "reading.docx")

    extraction.extract_course_text(course_dir=course_dir)

    assert "DOCX hello" in (course_dir / "Lecture 1" / "extracted_text.md").read_text(encoding="utf-8")


def test_pptx_extraction() -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "PPTX hello"
    presentation.save(course_dir / "Lecture 1" / "materials" / "slides.pptx")

    extraction.extract_course_text(course_dir=course_dir)

    assert "PPTX hello" in (course_dir / "Lecture 1" / "extracted_text.md").read_text(encoding="utf-8")


def test_html_extraction() -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    (course_dir / "Lecture 1" / "materials" / "page.html").write_text(
        "<html><body><main><h1>HTML hello</h1><script>skip</script></main></body></html>",
        encoding="utf-8",
    )

    extraction.extract_course_text(course_dir=course_dir)

    text = (course_dir / "Lecture 1" / "extracted_text.md").read_text(encoding="utf-8")
    assert "HTML hello" in text
    assert "skip" not in text


def test_txt_md_and_sql_code_extraction_as_plain_text() -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    materials = course_dir / "Lecture 1" / "materials"
    (materials / "notes.txt").write_text("TXT hello", encoding="utf-8")
    (materials / "notes.md").write_text("# MD hello", encoding="utf-8")
    (materials / "query.sql").write_text("select 'SQL hello';", encoding="utf-8")
    (materials / "script.py").write_text("print('PY hello')", encoding="utf-8")

    extraction.extract_course_text(course_dir=course_dir)

    text = (course_dir / "Lecture 1" / "extracted_text.md").read_text(encoding="utf-8")
    assert "TXT hello" in text
    assert "MD hello" in text
    assert "SQL hello" in text
    assert "PY hello" in text


def test_unsupported_file_skipping() -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    (course_dir / "Lecture 1" / "materials" / "image.png").write_bytes(b"png")

    summary = extraction.extract_course_text(course_dir=course_dir)

    manifest = read_section_manifest(course_dir)
    assert summary.files_skipped == 1
    assert manifest["skipped_unsupported_files"][0]["path"] == "materials/image.png"


def test_extraction_failure_is_recorded_without_crashing(monkeypatch: pytest.MonkeyPatch) -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    broken = course_dir / "Lecture 1" / "materials" / "broken.txt"
    broken.write_text("broken", encoding="utf-8")

    def fail_extract(path: Path):
        raise RuntimeError("boom")

    monkeypatch.setattr(extractors, "extract_file", fail_extract)

    summary = extraction.extract_course_text(course_dir=course_dir)

    manifest = read_section_manifest(course_dir)
    assert summary.files_failed == 1
    assert manifest["extraction_failures"][0]["path"] == "materials/broken.txt"
    assert "boom" in manifest["extraction_failures"][0]["error"]


def test_extracted_text_created_per_section_and_nested_materials_processed() -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    nested = course_dir / "Lecture 1" / "materials" / "Folder activity"
    nested.mkdir()
    (nested / "nested.txt").write_text("Nested hello", encoding="utf-8")

    extraction.extract_course_text(course_dir=course_dir)

    extracted = course_dir / "Lecture 1" / "extracted_text.md"
    text = extracted.read_text(encoding="utf-8")
    assert extracted.exists()
    assert "## Source: materials/Folder activity/nested.txt" in text
    assert "Nested hello" in text


def test_empty_text_is_recorded() -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    (course_dir / "Lecture 1" / "materials" / "empty.txt").write_text("", encoding="utf-8")

    extraction.extract_course_text(course_dir=course_dir)

    manifest = read_section_manifest(course_dir)
    text = (course_dir / "Lecture 1" / "extracted_text.md").read_text(encoding="utf-8")
    assert manifest["empty_text"][0]["path"] == "materials/empty.txt"
    assert "No extractable text found" in text


def test_manifest_updates() -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    (course_dir / "Lecture 1" / "materials" / "notes.txt").write_text("hello", encoding="utf-8")

    summary = extraction.extract_course_text(course_dir=course_dir)

    top_manifest = json.loads((course_dir / "manifest.json").read_text(encoding="utf-8"))
    section_manifest = read_section_manifest(course_dir)
    assert summary.sections_processed == 1
    assert top_manifest["extraction_summary"]["files_extracted"] == 1
    assert section_manifest["extracted_files"][0]["path"] == "materials/notes.txt"
    assert section_manifest["extracted_at"]


def test_cli_can_find_course_folder_by_course_id_prefix() -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    (course_dir / "Lecture 1" / "materials" / "notes.txt").write_text("CLI hello", encoding="utf-8")

    result = runner.invoke(app, ["text", "extract", "--course", "3025533", "--out", str(root)])

    assert result.exit_code == 0
    assert "Sections processed: 1" in result.output
    assert "CLI hello" in (course_dir / "Lecture 1" / "extracted_text.md").read_text(encoding="utf-8")


def test_cli_suppresses_noisy_pdf_warnings(monkeypatch: pytest.MonkeyPatch) -> None:
    root = local_tmp_path()
    course_dir = make_course(root)
    (course_dir / "Lecture 1" / "materials" / "lecture.pdf").write_bytes(b"%PDF noisy")

    def noisy_extract(path: Path):
        print("Ignoring wrong pointing object 34 0 (offset 0)", file=sys.stderr)
        return extractors.ExtractionResult(path=path, text="clean extracted text", extractor="pdf")

    monkeypatch.setattr(extractors, "extract_file", noisy_extract)

    result = runner.invoke(app, ["text", "extract", "--course-dir", str(course_dir)])

    assert result.exit_code == 0
    assert "Ignoring wrong pointing object" not in result.output
    assert "Files extracted: 1" in result.output
