from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def extract_text(path: str | Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"Slide {index}\n" + "\n".join(slide_text))
    return "\n\n".join(parts)
